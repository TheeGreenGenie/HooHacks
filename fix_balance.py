"""
fix_balance.py — Rework slide 7 from unit-economics breakdown
to a simple value-proposition slide. Keeps the visual card layout.
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

GOLD_LEAF = RGBColor(0xF5, 0xDF, 0x9A)
RUST      = RGBColor(0x8B, 0x25, 0x00)
BROWN     = RGBColor(0x4A, 0x2C, 0x14)
TAN       = RGBColor(0xB8, 0x96, 0x68)
F_H = "Rockwell"
F_B = "Georgia"

prs = Presentation("C:/Users/sgama/Downloads/HooHacks/Frontier_Finance_Pitch_v8.pptx")
slide7 = prs.slides[6]   # 0-indexed

# ── 1. Change heading ─────────────────────────────────────────────────────────
for shape in slide7.shapes:
    if shape.name == "TextBox 5":
        tf = shape.text_frame
        for para in tf.paragraphs:
            para.font.name = F_H
            for run in para.runs:
                run.text = "The Value Is Simple"
                run.font.name = F_H

# ── 2. Rework the 4 cards ─────────────────────────────────────────────────────
# New card content: (big line, small line)
new_cards = [
    ("$5 / mo",          "Flat monthly subscription"),
    ("$500 \u2013 $1,500",    "Estimated annual savings per user"),
    ("100\u00d7 return",      "On your subscription cost"),
    ("Paying $5\nto save hundreds\njust makes sense.", None),
]

card_idx = 0
for shape in slide7.shapes:
    if "Rounded Rectangle" not in shape.name:
        continue
    if card_idx >= len(new_cards):
        break
    big, small = new_cards[card_idx]
    tf = shape.text_frame
    tf.word_wrap = True

    # Clear all paragraph runs
    paras = tf.paragraphs
    # First para = big text
    p0 = paras[0]
    p0.alignment = PP_ALIGN.CENTER
    p0.font.name = F_H; p0.font.size = Pt(30 if card_idx < 3 else 18)
    p0.font.bold = True; p0.font.color.rgb = RUST
    for r in p0.runs: r.text = ""
    if p0.runs:
        p0.runs[0].text = big
        p0.runs[0].font.name = F_H
        p0.runs[0].font.size = Pt(30 if card_idx < 3 else 18)
        p0.runs[0].font.bold = True
        p0.runs[0].font.color.rgb = RUST
    else:
        from pptx.oxml.ns import qn
        from lxml import etree
        r_el = etree.SubElement(p0._p, qn('a:r'))
        rPr = etree.SubElement(r_el, qn('a:rPr'), attrib={'lang': 'en-US', 'dirty': '0'})
        t_el = etree.SubElement(r_el, qn('a:t'))
        t_el.text = big

    # Second para = small text
    if small and len(paras) > 1:
        p1 = paras[1]
        p1.alignment = PP_ALIGN.CENTER
        p1.font.name = F_B; p1.font.size = Pt(13)
        p1.font.bold = False; p1.font.color.rgb = BROWN
        for r in p1.runs: r.text = ""
        if p1.runs:
            p1.runs[0].text = small
            p1.runs[0].font.name = F_B
            p1.runs[0].font.size = Pt(13)
            p1.runs[0].font.bold = False
            p1.runs[0].font.color.rgb = BROWN
        else:
            from pptx.oxml.ns import qn
            from lxml import etree
            r_el = etree.SubElement(p1._p, qn('a:r'))
            rPr = etree.SubElement(r_el, qn('a:rPr'), attrib={'lang': 'en-US', 'dirty': '0'})
            t_el = etree.SubElement(r_el, qn('a:t'))
            t_el.text = small

    card_idx += 1

# ── 3. Remove the footnote / disclaimer from slide 7 (too jargony) ────────────
for shape in list(slide7.shapes):
    if shape.name in ("TextBox 10", "TextBox 11"):
        shape._element.getparent().remove(shape._element)

prs.save("C:/Users/sgama/Downloads/HooHacks/Frontier_Finance_Pitch_v8.pptx")
print("Done: slide 7 rebalanced.")
