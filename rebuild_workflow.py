"""
rebuild_workflow.py — Replace slide 5 logo grid with a proper workflow DAG.
Runs on top of the already-themed v8.pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# palette (matches theme_west.py)
GOLD      = RGBColor(0xC4, 0x89, 0x1A)
DARK_GOLD = RGBColor(0x9A, 0x6D, 0x07)
GOLD_LEAF = RGBColor(0xF5, 0xDF, 0x9A)
TAN       = RGBColor(0xB8, 0x96, 0x68)
LEATHER   = RGBColor(0x2A, 0x13, 0x04)

C_INFRA    = RGBColor(0x2A, 0x13, 0x04)
C_USER     = RGBColor(0x1E, 0x0E, 0x03)
C_OCR      = RGBColor(0x42, 0x26, 0x0A)
C_AI       = RGBColor(0x54, 0x32, 0x10)
C_STOCK    = RGBColor(0x38, 0x20, 0x0B)
C_VOICE    = RGBColor(0x60, 0x3A, 0x14)
C_PAY      = RGBColor(0x70, 0x48, 0x18)
C_FALLBACK = RGBColor(0x7A, 0x58, 0x30)
C_DB       = RGBColor(0x8B, 0x25, 0x00)

F_H = "Rockwell"
F_B = "Georgia"

# ── geometry ──────────────────────────────────────────────────────────────────
ROW_A_Y = 1.72
ROW_A_H = 0.44

# auth chain row  (User, Auth0, Next.js, FastAPI)
USER_X, USER_W = 0.28, 1.22
AUTH_X, AUTH_W = 1.72, 1.28
NEXT_X, NEXT_W = 3.20, 2.08
FAPI_X, FAPI_W = 5.52, 2.70   # center = 6.87"
FAPI_CX = FAPI_X + FAPI_W / 2  # 6.87

# 5 columns
BWIDTH = 2.46
BGAP   = 0.10
B_X  = [0.22 + i * (BWIDTH + BGAP) for i in range(5)]
B_CX = [x + BWIDTH / 2 for x in B_X]
# B_CX = [1.45, 4.01, 6.57, 9.13, 11.69]

STEM_Y0 = ROW_A_Y + ROW_A_H       # 2.16
BUS_Y   = STEM_Y0 + 0.12          # 2.28
BANNER_Y = BUS_Y + 0.12           # 2.40
BANNER_H = 0.27

BOX_Y0   = BANNER_Y + BANNER_H + 0.11  # 2.78
BOX_H    = 0.50
BOX_VGAP = 0.10

def box_y(r): return BOX_Y0 + r * (BOX_H + BOX_VGAP)
# 0->2.78  1->3.38  2->3.98  3->4.58

MONGO_Y = box_y(3) + BOX_H + 0.20   # 5.28
MONGO_H = 0.58


# ── helpers ───────────────────────────────────────────────────────────────────
def node(slide, x, y, w, h, fill, text, size=10.5, bold=True,
         color=GOLD_LEAF, rounded=True, bdr=None, bdr_pt=1.5):
    sh = slide.shapes.add_shape(
        5 if rounded else 1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = bdr if bdr else GOLD
    sh.line.width = Pt(bdr_pt)
    tf = sh.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after  = Pt(0)
        p.font.name = F_H; p.font.size = Pt(size)
        p.font.bold = bold; p.font.color.rgb = color
        r = p.add_run()
        r.text = ln
        r.font.name = F_H; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return sh


def banner(slide, x, y, w, h, fill, text):
    sh = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = DARK_GOLD
    sh.line.width = Pt(0.75)
    tf = sh.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = F_H; p.font.size = Pt(9); p.font.bold = True
    p.font.color.rgb = GOLD_LEAF
    r = p.add_run()
    r.text = text
    r.font.name = F_H; r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = GOLD_LEAF
    return sh


def hline(slide, x0, y, x1, color=GOLD, pt=1.5):
    h_emu = max(int(Pt(pt)), 9525)
    sh = slide.shapes.add_shape(
        1, Inches(x0), int(Inches(y)) - h_emu // 2,
        Inches(x1 - x0), h_emu)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def vline(slide, x, y0, y1, color=GOLD, pt=1.5):
    w_emu = max(int(Pt(pt)), 9525)
    sh = slide.shapes.add_shape(
        1, int(Inches(x)) - w_emu // 2, Inches(y0),
        w_emu, Inches(y1 - y0))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def tiny_label(slide, x, y, w, h, text, color=TAN, size=7.5):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = F_B; p.font.size = Pt(size)
    p.font.italic = True; p.font.color.rgb = color
    r = p.add_run()
    r.text = text
    r.font.name = F_B; r.font.size = Pt(size)
    r.font.italic = True; r.font.color.rgb = color


# ── load & clean slide 5 ─────────────────────────────────────────────────────
prs  = Presentation("C:/Users/sgama/Downloads/HooHacks/Frontier_Finance_Pitch_v8.pptx")
slide = prs.slides[4]

KEEP = {"Rectangle 1", "Rectangle 2", "Rectangle 3", "Diamond 4", "TextBox 5"}
dead = [s for s in slide.shapes if s.name not in KEEP]
for s in dead:
    s._element.getparent().remove(s._element)

# re-add gold trim lines (removed above since they had auto-names)
trim_top = slide.shapes.add_shape(1,
    Inches(0), int(Inches(0.38)), int(Inches(13.33)), int(Pt(1.5)))
trim_top.fill.solid(); trim_top.fill.fore_color.rgb = GOLD
trim_top.line.fill.background()

trim_bot = slide.shapes.add_shape(1,
    Inches(0), int(Inches(7.10)), int(Inches(13.33)), int(Pt(1.5)))
trim_bot.fill.solid(); trim_bot.fill.fore_color.rgb = GOLD
trim_bot.line.fill.background()

# ── ROW A: auth chain ─────────────────────────────────────────────────────────
node(slide, USER_X, ROW_A_Y, USER_W, ROW_A_H, C_USER,  "User\nBrowser",      size=9.5)
node(slide, AUTH_X, ROW_A_Y, AUTH_W, ROW_A_H, C_INFRA, "Auth0\n+ TOTP",      size=9.5)
node(slide, NEXT_X, ROW_A_Y, NEXT_W, ROW_A_H, C_INFRA, "Next.js 15\nFrontend", size=9.5)
node(slide, FAPI_X, ROW_A_Y, FAPI_W, ROW_A_H, C_INFRA, "FastAPI\nBackend",   size=9.5)

# horizontal connectors in row A
mid_a = ROW_A_Y + ROW_A_H / 2
hline(slide, USER_X + USER_W, mid_a, AUTH_X)
hline(slide, AUTH_X + AUTH_W, mid_a, NEXT_X)
hline(slide, NEXT_X + NEXT_W, mid_a, FAPI_X)

# ── bus topology: FastAPI fans out to 5 columns ───────────────────────────────
vline(slide, FAPI_CX, STEM_Y0, BUS_Y)                    # stem down
hline(slide, B_CX[0], BUS_Y, B_CX[4])                    # horizontal bus
for cx in B_CX:
    vline(slide, cx, BUS_Y, BANNER_Y)                     # drops to banners

# ── category banners ──────────────────────────────────────────────────────────
banner(slide, B_X[0], BANNER_Y, BWIDTH, BANNER_H, C_OCR,   "OCR Pipeline")
banner(slide, B_X[1], BANNER_Y, BWIDTH, BANNER_H, C_AI,    "AI Suggestions")
banner(slide, B_X[2], BANNER_Y, BWIDTH, BANNER_H, C_STOCK, "Stock Trading")
banner(slide, B_X[3], BANNER_Y, BWIDTH, BANNER_H, C_VOICE, "Voice Chat")
banner(slide, B_X[4], BANNER_Y, BWIDTH, BANNER_H, C_PAY,   "Payments")

# ── column 0: OCR Pipeline ───────────────────────────────────────────────────
node(slide, B_X[0], box_y(0), BWIDTH, BOX_H, C_OCR, "Google\nCloud Vision")
vline(slide, B_CX[0], box_y(0) + BOX_H, box_y(1))
node(slide, B_X[0], box_y(1), BWIDTH, BOX_H, C_OCR, "Transaction\nParser")
vline(slide, B_CX[0], box_y(1) + BOX_H, MONGO_Y)

# ── column 1: AI Suggestions (4-level fallback chain) ────────────────────────
node(slide, B_X[1], box_y(0), BWIDTH, BOX_H, C_AI,     "Gemini")
vline(slide, B_CX[1], box_y(0) + BOX_H, box_y(1))
tiny_label(slide, B_CX[1] + 0.06, box_y(0) + BOX_H + 0.01, 0.80, 0.10, "fallback")
node(slide, B_X[1], box_y(1), BWIDTH, BOX_H, C_AI,     "Snowflake\nCortex")
vline(slide, B_CX[1], box_y(1) + BOX_H, box_y(2))
tiny_label(slide, B_CX[1] + 0.06, box_y(1) + BOX_H + 0.01, 0.80, 0.10, "fallback")
node(slide, B_X[1], box_y(2), BWIDTH, BOX_H, C_AI,     "OpenRouter")
vline(slide, B_CX[1], box_y(2) + BOX_H, box_y(3))
tiny_label(slide, B_CX[1] + 0.06, box_y(2) + BOX_H + 0.01, 0.80, 0.10, "fallback")
node(slide, B_X[1], box_y(3), BWIDTH, BOX_H, C_FALLBACK,"Local\nFallback", bdr=DARK_GOLD)
vline(slide, B_CX[1], box_y(3) + BOX_H, MONGO_Y, color=TAN, pt=1.0)  # reads MongoDB

# ── column 2: Stock Trading ───────────────────────────────────────────────────
node(slide, B_X[2], box_y(0), BWIDTH, BOX_H, C_STOCK, "yfinance /\nyahooquery")
vline(slide, B_CX[2], box_y(0) + BOX_H, box_y(1))
node(slide, B_X[2], box_y(1), BWIDTH, BOX_H, C_STOCK, "TF LSTM\n+ sklearn")
vline(slide, B_CX[2], box_y(1) + BOX_H, box_y(2))
node(slide, B_X[2], box_y(2), BWIDTH, BOX_H, C_STOCK, "Plotly\nCharts")
vline(slide, B_CX[2], box_y(2) + BOX_H, MONGO_Y)

# ── column 3: Voice Chat ──────────────────────────────────────────────────────
node(slide, B_X[3], box_y(0), BWIDTH, BOX_H, C_VOICE, "Gemini\nNLP")   # Gemini again - OK per user
vline(slide, B_CX[3], box_y(0) + BOX_H, box_y(1))
node(slide, B_X[3], box_y(1), BWIDTH, BOX_H, C_VOICE,    "ElevenLabs\nTTS")
vline(slide, B_CX[3], box_y(1) + BOX_H, box_y(2))
node(slide, B_X[3], box_y(2), BWIDTH, BOX_H, C_FALLBACK, "Browser\nSpeech API", bdr=DARK_GOLD)
tiny_label(slide, B_X[3], box_y(2) - 0.12, BWIDTH, 0.12, "fallback", size=7)
# voice is stateless — no line to MongoDB

# ── column 4: Payments ────────────────────────────────────────────────────────
node(slide, B_X[4], box_y(0), BWIDTH, BOX_H, C_PAY, "Stripe\n(fiat)")
vline(slide, B_CX[4], box_y(0) + BOX_H, box_y(1))
node(slide, B_X[4], box_y(1), BWIDTH, BOX_H, C_PAY, "Solana\n(crypto)")
vline(slide, B_CX[4], box_y(1) + BOX_H, box_y(2))
node(slide, B_X[4], box_y(2), BWIDTH, BOX_H, C_PAY, "Wolfram\nAlpha")
vline(slide, B_CX[4], box_y(2) + BOX_H, MONGO_Y)

# ── MongoDB bar (full width, spans all columns) ───────────────────────────────
MONGO_X = B_X[0]
MONGO_W = B_X[4] + BWIDTH - B_X[0]
mongo = slide.shapes.add_shape(
    1, Inches(MONGO_X), Inches(MONGO_Y), Inches(MONGO_W), Inches(MONGO_H))
mongo.fill.solid(); mongo.fill.fore_color.rgb = C_DB
mongo.line.color.rgb = GOLD; mongo.line.width = Pt(2.0)
tf = mongo.text_frame
p  = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.name = F_H; p.font.size = Pt(13); p.font.bold = True
p.font.color.rgb = GOLD_LEAF
r = p.add_run()
r.text = "MongoDB  —  users  \u00b7  transactions  \u00b7  portfolios  \u00b7  stock positions"
r.font.name = F_H; r.font.size = Pt(13); r.font.bold = True
r.font.color.rgb = GOLD_LEAF

prs.save("C:/Users/sgama/Downloads/HooHacks/Frontier_Finance_Pitch_v8.pptx")
print("Done: workflow rebuilt in v8.pptx")
