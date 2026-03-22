"""
theme_west.py — Apply Wild West theme to Frontier Finance pitch deck.
Output: Frontier_Finance_Pitch_v8.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
PARCHMENT   = RGBColor(0xEF, 0xE2, 0xCA)  # aged parchment (slide bg)
LEATHER     = RGBColor(0x2A, 0x13, 0x04)  # dark saddle leather (top/bot bars)
GOLD        = RGBColor(0xC4, 0x89, 0x1A)  # burnished gold (diamond accent)
RUST        = RGBColor(0x8B, 0x25, 0x00)  # burnt sienna (stat numbers, highlights)
CREAM       = RGBColor(0xFA, 0xF3, 0xE2)  # bleached cream (card bg)
DARK_GOLD   = RGBColor(0x9A, 0x6D, 0x07)  # dark gold (card borders, rule lines)
INK         = RGBColor(0x18, 0x0C, 0x02)  # near-black ink (headings)
BROWN       = RGBColor(0x4A, 0x2C, 0x14)  # medium brown (body text)
TAN         = RGBColor(0xB8, 0x96, 0x68)  # warm tan (dividers, footnotes)
GOLD_LEAF   = RGBColor(0xF5, 0xDF, 0x9A)  # gold leaf (text on dark bg)
WHEAT       = RGBColor(0xE2, 0xCB, 0x96)  # wheat (subtitle on dark bg)
DARK_BG     = RGBColor(0x16, 0x07, 0x00)  # near-black sepia (title overlay)
DIVIDER     = RGBColor(0xC0, 0x9A, 0x68)  # vertical divider line

# ── Fonts ─────────────────────────────────────────────────────────────────────
F_DISPLAY = "Rockwell Extra Bold"   # hero title slides
F_HEADING = "Rockwell"              # section headings
F_BODY    = "Georgia"               # body / bullet text
F_STAT    = "Rockwell Extra Bold"   # big numbers
F_LABEL   = "Rockwell"              # card titles, tech labels


# ── Helpers ───────────────────────────────────────────────────────────────────

def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def border(shape, color, pt=2.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(pt)

def no_border(shape):
    try:
        shape.line.fill.background()
    except Exception:
        pass

def set_fill_alpha(shape, color_rgb, alpha_pct):
    """Solid fill with % opacity (100=opaque, 0=transparent)."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color_rgb
    xFill = shape.fill._xPr
    solid_el = xFill.find(qn('a:solidFill'))
    if solid_el is None:
        return
    clr_el = solid_el.find(qn('a:srgbClr'))
    if clr_el is None:
        return
    for a in clr_el.findall(qn('a:alpha')):
        clr_el.remove(a)
    alpha_el = etree.SubElement(clr_el, qn('a:alpha'))
    alpha_el.set('val', str(int(alpha_pct * 1000)))

def style_para(para, font=None, size=None, bold=None,
               italic=None, color=None, align=None):
    if align is not None:
        para.alignment = align
    pf = para.font
    if font:   pf.name = font
    if size:   pf.size = Pt(size)
    if bold is not None:   pf.bold = bold
    if italic is not None: pf.italic = italic
    if color:  pf.color.rgb = color
    for run in para.runs:
        if font:   run.font.name = font
        if size:   run.font.size = Pt(size)
        if bold is not None:   run.font.bold = bold
        if italic is not None: run.font.italic = italic
        if color:  run.font.color.rgb = color

def style_tf(tf, font=None, size=None, bold=None,
             italic=None, color=None, align=None):
    for para in tf.paragraphs:
        style_para(para, font, size, bold, italic, color, align)

def add_rule(slide, y_in, color, h_pt=1.8, x0=0.58, x1=12.75):
    """Thin horizontal decorative rule."""
    sh = slide.shapes.add_shape(
        1,  # rectangle
        Inches(x0), Inches(y_in),
        Inches(x1 - x0), Pt(h_pt)
    )
    solid(sh, color)
    no_border(sh)
    return sh

def add_double_rule(slide, y_in, color):
    """Two thin rules close together for a western frame feel."""
    add_rule(slide, y_in,         color, h_pt=1.5)
    add_rule(slide, y_in + 0.04,  color, h_pt=0.8)


# ── Main ──────────────────────────────────────────────────────────────────────
prs = Presentation("C:/Users/sgama/Downloads/HooHacks/Frontier_Finance_Pitch_v7.pptx")
SW = prs.slide_width
SH = prs.slide_height

for si, slide in enumerate(prs.slides):
    sn = si + 1
    sh = {s.name: s for s in slide.shapes}

    # =========================================================================
    # SLIDE 1 — Hero title
    # =========================================================================
    if sn == 1:
        if "Rectangle 2" in sh:
            set_fill_alpha(sh["Rectangle 2"], DARK_BG, 78)   # let bg image breathe
        if "TextBox 3" in sh:
            style_tf(sh["TextBox 3"].text_frame, F_DISPLAY, 66, bold=True, color=GOLD_LEAF)
        if "TextBox 4" in sh:
            style_tf(sh["TextBox 4"].text_frame, F_BODY, 22, italic=True, color=WHEAT)

    # =========================================================================
    # SLIDES 2–12 — Content slides
    # =========================================================================
    elif 2 <= sn <= 12:

        # Background
        if "Rectangle 1" in sh:
            solid(sh["Rectangle 1"], PARCHMENT)

        # Top bar — slightly taller, with a thin gold trim line inside
        if "Rectangle 2" in sh:
            s = sh["Rectangle 2"]
            solid(s, LEATHER)
            s.height = Inches(0.40)

        # Bottom bar
        if "Rectangle 3" in sh:
            s = sh["Rectangle 3"]
            solid(s, LEATHER)
            s.height = Inches(0.40)
            s.top = int(SH - Inches(0.40))

        # Gold trim line just inside the top bar
        trim = slide.shapes.add_shape(1,
            Inches(0), Inches(0.38),
            SW, Pt(1.5))
        solid(trim, GOLD)
        no_border(trim)

        # Gold trim line just inside the bottom bar
        trim2 = slide.shapes.add_shape(1,
            Inches(0), int(SH - Inches(0.40)),
            SW, Pt(1.5))
        solid(trim2, GOLD)
        no_border(trim2)

        # Diamond accent — larger, richer gold
        if "Diamond 4" in sh:
            d = sh["Diamond 4"]
            solid(d, GOLD)
            d.width  = int(Inches(0.24))
            d.height = int(Inches(0.24))
            d.left   = int(SW / 2 - Inches(0.12))
            d.top    = int(Inches(0.08))

        # Section heading
        if "TextBox 5" in sh:
            style_tf(sh["TextBox 5"].text_frame, F_HEADING, 33, bold=True, color=INK)

        # Decorative double rule under heading
        add_double_rule(slide, 1.54, DARK_GOLD)

        # Vertical divider (slides 2, 8)
        if "Rectangle 6" in sh:
            solid(sh["Rectangle 6"], DIVIDER)
            sh["Rectangle 6"].width = Pt(1.5)

        # ── Rounded Rectangle cards (slides 4, 7, 12) ────────────────────────
        for shape in slide.shapes:
            if "Rounded Rectangle" not in shape.name:
                continue
            solid(shape, CREAM)
            border(shape, DARK_GOLD, 2.5)
            tf = shape.text_frame
            tf.word_wrap = True
            paras = [p for p in tf.paragraphs if p.text.strip()]
            for pi, para in enumerate(paras):
                if pi == 0:
                    style_para(para, F_STAT, 34, bold=True, color=RUST,
                               align=PP_ALIGN.CENTER)
                else:
                    style_para(para, F_BODY, 14, bold=False, color=BROWN,
                               align=PP_ALIGN.CENTER)

        # ── Slide-specific text ───────────────────────────────────────────────

        # Slide 2 — "The money gap is real"
        if sn == 2:
            # Left bullets
            if "TextBox 7" in sh:
                tf = sh["TextBox 7"].text_frame
                for para in tf.paragraphs:
                    style_para(para, F_BODY, 17, color=BROWN)
            # Right stat "4.2%"
            if "TextBox 8" in sh:
                paras = sh["TextBox 8"].text_frame.paragraphs
                for pi, para in enumerate(paras):
                    if pi == 0:
                        style_para(para, F_STAT, 72, bold=True, color=RUST)
                    else:
                        style_para(para, F_BODY, 17, color=BROWN)

        # Slide 3 — "Financial literacy remains low"
        elif sn == 3:
            for tname in ["TextBox 6", "TextBox 7"]:
                if tname not in sh: continue
                paras = sh[tname].text_frame.paragraphs
                for pi, para in enumerate(paras):
                    if pi == 0:
                        style_para(para, F_STAT, 72, bold=True, color=RUST)
                    else:
                        style_para(para, F_BODY, 17, color=BROWN)

        # Slide 4 — cards are handled above by Rounded Rectangle loop

        # Slide 5 — tech workflow (logo labels)
        elif sn == 5:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.name != "TextBox 5":
                    style_tf(shape.text_frame, F_LABEL, 13, bold=True, color=BROWN,
                             align=PP_ALIGN.CENTER)

        # Slide 6 — Savings Impact
        elif sn == 6:
            if "TextBox 6" in sh:
                paras = sh["TextBox 6"].text_frame.paragraphs
                for pi, para in enumerate(paras):
                    if pi == 0:
                        style_para(para, F_STAT, 72, bold=True, color=RUST)
                    else:
                        style_para(para, F_BODY, 17, color=BROWN)
            if "TextBox 7" in sh:
                style_tf(sh["TextBox 7"].text_frame, F_BODY, 18, italic=True, color=BROWN)

        # Slide 7 — Business Model (cards handled above)

        # Slide 8 — Equity + Accessibility
        elif sn == 8:
            if "TextBox 7" in sh:
                style_tf(sh["TextBox 7"].text_frame, F_BODY, 17, color=BROWN)
            if "TextBox 8" in sh:
                style_tf(sh["TextBox 8"].text_frame, F_BODY, 17, color=BROWN)

        # Slide 9 — Built Today
        elif sn == 9:
            if "TextBox 6" in sh:
                style_tf(sh["TextBox 6"].text_frame, F_BODY, 19, color=INK)

        # Slide 10 — Demo video (heading already styled above)

        # Slide 11 — Tech stack logos (labels)
        elif sn == 11:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.name != "TextBox 5":
                    style_tf(shape.text_frame, F_LABEL, 13, bold=True, color=BROWN,
                             align=PP_ALIGN.CENTER)

        # Slide 12 — Next 30-60 Days (cards handled above)

        # Slide 7 extra: disclaimer + footnotes already caught below

        # ── Footnotes — anything in bottom 0.65" of slide ────────────────────
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.name == "TextBox 5":
                continue
            top_in = shape.top / 914400
            if top_in > 6.85:
                style_tf(shape.text_frame, F_BODY, 9.5, italic=True, color=TAN)

        # Slide 7 disclaimer (TextBox 10 at y≈6.20")
        if sn == 7 and "TextBox 10" in sh:
            style_tf(sh["TextBox 10"].text_frame, F_BODY, 10, italic=True, color=TAN)

    # =========================================================================
    # SLIDE 13 — Outro
    # =========================================================================
    elif sn == 13:
        if "Rectangle 2" in sh:
            set_fill_alpha(sh["Rectangle 2"], DARK_BG, 78)
        if "TextBox 3" in sh:
            style_tf(sh["TextBox 3"].text_frame, F_DISPLAY, 66, bold=True, color=GOLD_LEAF)
        if "TextBox 4" in sh:
            style_tf(sh["TextBox 4"].text_frame, F_BODY, 22, italic=True, color=WHEAT)


prs.save("C:/Users/sgama/Downloads/HooHacks/Frontier_Finance_Pitch_v8.pptx")
print("Saved: Frontier_Finance_Pitch_v8.pptx")
